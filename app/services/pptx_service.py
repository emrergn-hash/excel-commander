"""
Excel Commander - PowerPoint Generation Service
Creates professional PPTX files from Excel data.
"""
import os
import uuid
import logging
from typing import List, Any, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

logger = logging.getLogger(__name__)

# Output directory for generated files
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "generated")


class PPTXService:
    """Service class for PowerPoint generation."""
    
    # Excel Commander Brand Colors
    COLOR_PRIMARY = RGBColor(16, 124, 16)      # Excel Green #107C10
    COLOR_SECONDARY = RGBColor(0, 120, 212)    # Office Blue #0078D4
    COLOR_DARK = RGBColor(50, 50, 50)          # Dark Gray
    COLOR_LIGHT = RGBColor(243, 242, 241)      # Light Gray
    
    def __init__(self):
        # Ensure output directory exists
        os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    def create_presentation(
        self,
        data: List[List[Any]],
        title: str = "Analiz Raporu",
        insights: Optional[List[str]] = None,
        include_chart: bool = True,
        chart_type: str = "bar"
    ) -> str:
        """
        Create a professional PowerPoint presentation from Excel data.
        
        Args:
            data: 2D array with headers in first row
            title: Presentation title
            insights: List of AI-generated insights
            include_chart: Whether to include a chart slide
            chart_type: Type of chart (bar, line, pie)
        
        Returns:
            Path to generated PPTX file
        """
        prs = Presentation()
        
        # Set slide dimensions (16:9 widescreen)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title Slide
        self._add_title_slide(prs, title)
        
        # Slide 2: Key Insights
        if insights:
            self._add_insights_slide(prs, insights)
        
        # Slide 3: Data Chart
        if include_chart and len(data) > 1:
            self._add_chart_slide(prs, data, chart_type)
        
        # Slide 4: Data Table
        if len(data) > 1:
            self._add_table_slide(prs, data)
        
        # Slide 5: Conclusion
        self._add_conclusion_slide(prs, title)
        
        # Save file
        filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        
        logger.info(f"Presentation created: {filepath}")
        return filepath
    
    def _add_title_slide(self, prs: Presentation, title: str):
        """Add a stylish title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Background shape (green bar)
        left = Inches(0)
        top = Inches(2.5)
        width = prs.slide_width
        height = Inches(2.5)
        
        shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLOR_PRIMARY
        shape.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.5))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Excel Commander ile Oluşturuldu • {datetime.now().strftime('%d.%m.%Y')}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_insights_slide(self, prs: Presentation, insights: List[str]):
        """Add a slide with key insights."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "📊 Önemli Bulgular"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_DARK
        
        # Insights as bullet points with icons
        for i, insight in enumerate(insights[:5]):
            y_pos = Inches(1.8 + i * 1.1)
            
            # Bullet box
            box = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11), Inches(0.9))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = insight
            p.font.size = Pt(20)
            p.font.color.rgb = self.COLOR_DARK
            
            # Add left border indicator
            indicator = slide.shapes.add_shape(1, Inches(0.5), y_pos, Inches(0.15), Inches(0.8))
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = self.COLOR_PRIMARY
            indicator.line.fill.background()
    
    def _add_chart_slide(self, prs: Presentation, data: List[List[Any]], chart_type: str):
        """Add a slide with a data chart."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "📈 Veri Analizi"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_DARK
        
        # Prepare chart data
        headers = data[0]
        rows = data[1:]
        
        chart_data = CategoryChartData()
        
        # Categories (first column values)
        categories = [str(row[0]) for row in rows]
        chart_data.categories = categories
        
        # Series (remaining columns)
        for col_idx in range(1, len(headers)):
            series_name = str(headers[col_idx])
            series_values = []
            for row in rows:
                try:
                    val = float(row[col_idx]) if col_idx < len(row) else 0
                except (ValueError, TypeError):
                    val = 0
                series_values.append(val)
            chart_data.add_series(series_name, series_values)
        
        # Determine chart type
        if chart_type == "line":
            xl_chart_type = XL_CHART_TYPE.LINE
        elif chart_type == "pie":
            xl_chart_type = XL_CHART_TYPE.PIE
        else:
            xl_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        
        # Add chart
        x, y, cx, cy = Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8)
        chart = slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, chart_data).chart
        
        # Style the chart
        chart.has_legend = True
        chart.legend.include_in_layout = False
    
    def _add_table_slide(self, prs: Presentation, data: List[List[Any]]):
        """Add a slide with a data table."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "📋 Veri Tablosu"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_DARK
        
        # Limit rows for readability
        display_data = data[:15]  # Max 15 rows
        rows_count = len(display_data)
        cols_count = min(len(display_data[0]), 6)  # Max 6 columns
        
        # Calculate table dimensions
        table_width = Inches(12)
        table_height = Inches(rows_count * 0.45)
        
        # Add table
        table = slide.shapes.add_table(
            rows_count, cols_count,
            Inches(0.6), Inches(1.3),
            table_width, table_height
        ).table
        
        # Style header row
        for col_idx in range(cols_count):
            cell = table.cell(0, col_idx)
            cell.text = str(display_data[0][col_idx]) if col_idx < len(display_data[0]) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLOR_PRIMARY
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Style data rows
        for row_idx in range(1, rows_count):
            for col_idx in range(cols_count):
                cell = table.cell(row_idx, col_idx)
                value = display_data[row_idx][col_idx] if col_idx < len(display_data[row_idx]) else ""
                cell.text = str(value)
                
                # Alternating row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLOR_LIGHT
                
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = self.COLOR_DARK
                paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_conclusion_slide(self, prs: Presentation, title: str):
        """Add a conclusion/thank you slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.COLOR_PRIMARY
        bg_shape.line.fill.background()
        
        # Thank you text
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Teşekkürler"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"🚀 Excel Commander ile hazırlandı"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.alignment = PP_ALIGN.CENTER


# Singleton instance
_pptx_service: Optional[PPTXService] = None

def get_pptx_service() -> PPTXService:
    """Get or create PPTXService singleton."""
    global _pptx_service
    if _pptx_service is None:
        _pptx_service = PPTXService()
    return _pptx_service
